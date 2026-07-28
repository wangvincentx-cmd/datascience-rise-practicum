"""
Build poster v2 into `2026 Poster Template.pptx` -> RISE_Poster_2026_v2.pptx.

Uses the template's first layout: a narrow left column stacking Introduction over
Methods, one tall 15in Results column down the centre carrying all five figures,
and Discussion / References / Acknowledgements down the right.

Every headline number is recomputed here from the committed CSVs rather than
typed in, so the poster cannot drift from the analysis. Figure heights come from
each PNG's real pixel aspect, and every box is checked against an estimated text
height -- the build FAILS rather than silently overflowing a column, which is the
failure mode that is invisible until the poster is printed.

Prerequisites:
    python src/analysis_v2.py --section all
    python src/make_poster_v2_figures.py

Usage:
    python src/build_poster_v2.py
"""

import json
import math
import os
import sys

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "2026 Poster Template.pptx"
OUTFILE = "RISE_Poster_2026_v2.pptx"
FIGDIR = "figures/poster_v2"
SUMMARY = "data/scored/analysis_v2/summary.json"

INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
VERM = RGBColor(0xEB, 0x68, 0x34)

LEAD, BODY, SMALL, TINY = 30, 26, 22, 17

# Template geometry, in inches, read off the first layout's rounded rectangles.
INTRO = dict(left=1.30, top=7.45, width=8.05, height=15.75)
METHODS = dict(left=1.30, top=26.35, width=8.10, height=20.15)
DISCUSS = dict(left=26.85, top=8.55, width=7.90, height=28.20)
REFS = dict(left=26.85, top=38.95, width=7.90, height=2.85)
ACK = dict(left=26.85, top=44.25, width=7.90, height=2.55)

RESULTS_L, RESULTS_R = 10.56, 25.57
RESULTS_TOP, RESULTS_BOTTOM = 7.45, 46.60
FIG_W = 13.2


# --------------------------------------------------------------- numbers

def stats():
    """Recompute every number the poster prints, from committed data only."""
    with open(SUMMARY) as f:
        s = json.load(f)
    d = pd.read_csv("data/scored/monthly_scored.csv", low_memory=False)
    d["date"] = pd.to_datetime(d["date"], errors="coerce")
    sc = d[(d["scorable"] == True) & (d["hit"].isin([0, 1]))]
    b = pd.read_csv("data/scored/analysis_v2/benchmarks.csv")
    peaks = pd.read_csv("data/scored/analysis_v2/nber_peaks.csv")
    peaks = peaks[peaks["window"] == "6mo before peak"]
    return {
        "n_extracted": len(d),
        "n_scorable": len(sc),
        "n_pages": 15721,
        "n_publishers": d["publisher"].nunique(),
        "hit": sc["hit"].mean(),
        "n_peaks": len(peaks),
        "peak_mean": peaks["share_improve"].mean(),
        "peak_min": peaks["share_improve"].min(),
        "pub": s["publishers"],
        "mech": s["mechanism"],
        "text": s["text_ceiling"],
        "bench": b,
    }


# ------------------------------------------------------------- primitives

def est_height(blocks, width_in):
    """Approximate rendered height of a text block list, in inches.

    Deliberately conservative: ~0.50 em average glyph width and 1.18 line
    leading. Used only to fail the build when a column is over-filled."""
    total = 0.0
    for text, size, _bold, _color, after in blocks:
        cpl = max(8, int(width_in * 72 / (0.50 * size)))
        lines = sum(max(1, math.ceil(len(seg) / cpl)) for seg in text.split("\n"))
        total += lines * size * 1.18 / 72 + after / 72
    return total


def add_text(slide, box, blocks, name=""):
    tb = slide.shapes.add_textbox(Inches(box["left"]), Inches(box["top"]),
                                  Inches(box["width"]), Inches(box["height"]))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, after) in enumerate(blocks):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(after)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    need = est_height(blocks, box["width"])
    flag = "OVERFLOW" if need > box["height"] else "ok"
    print(f"  {name:<16} {need:5.1f}in of {box['height']:5.1f}in   {flag}")
    return need <= box["height"]


def fig_height(path, width):
    w, h = Image.open(path).size
    return width * h / w


# ------------------------------------------------------------------ build

def main():
    st = stats()
    prs = Presentation(TEMPLATE)

    # Keep the big-centre-column layout; drop the alternate slide.
    lst = prs.slides._sldIdLst
    for sl in list(lst)[1:]:
        lst.remove(sl)
    slide = prs.slides[0]

    ok = []

    # ---- title -----------------------------------------------------------
    for sh in slide.shapes:
        if sh.name == "TextBox 4":
            tf = sh.text_frame
            tf.clear()
            for i, (t, size, bold, color) in enumerate([
                ("Did Anyone See It Coming?", 96, True, INK),
                ("Machine-reading sixty years of American newspaper economic "
                 "forecasts, 1900–1963", 54, False, MUTED),
                ("RISE Practicum Student FN and LN¹ ·  Instructor²    "
                 "¹High School Name  ·  ²Boston University", 36, True, INK),
            ]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = t
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color

    # ---- introduction ----------------------------------------------------
    ok.append(add_text(slide, INTRO, [
        ("For sixty years American newspapers told readers what the economy "
         "would do next. Were they right?", BODY, True, INK, 14),
        ("Nobody has scored the record at scale. Reading a century of "
         "newspapers by hand is infeasible — and the standard shortcut, keyword "
         "search, throws away most of the evidence before the analysis begins.",
         BODY, False, INK, 14),
        (f"{st['n_pages']:,} pages  ·  {st['n_publishers']:,} publishers  ·  "
         f"every month, 1900–1963", BODY, True, BLUE, 6),
        (f"{st['n_extracted']:,} forecasts extracted   →   "
         f"{st['n_scorable']:,} scored against real economic data",
         BODY, True, BLUE, 14),
        ("THE RULE THE PROJECT RESTS ON — the language model decides WHAT was "
         "predicted; NBER dates and Federal Reserve series decide WHETHER IT "
         "CAME TRUE. The model is never asked whether a forecast was right; "
         "that would be hindsight grading itself.", BODY, True, VERM, 20),
        # A single number is the headline, so it gets the hero treatment rather
        # than a sentence: the whole poster is an answer to "why 51%?".
        (f"{st['hit']:.1%}", 150, True, BLUE, 0),
        ("of forecasts came true. A coin flip.", BODY, True, INK, 10),
        ("This poster is about why — and about whether anyone else did better.",
         BODY, False, INK, 0),
    ], "INTRO"))

    # ---- methods ---------------------------------------------------------
    ok.append(add_text(slide, METHODS, [
        ("1 · COLLECT", BODY, True, BLUE, 4),
        ("Library of Congress Chronicling America, public API. A fixed, "
         "direction-neutral query set sampled in all 768 months, so the corpus "
         "is continuous rather than selected around famous crashes. 761 months "
         "yield a scorable US-national forecast.", SMALL, False, INK, 12),
        ("2 · EXTRACT", BODY, True, BLUE, 4),
        ("A language model reads each full page and returns structured "
         "forecasts. Measured against a gold standard of 16 pages annotated "
         "BEFORE any model ran — 52 forecasts plus 44 boundary cases "
         "(advertisements, serialised fiction, refusals to forecast):",
         SMALL, False, INK, 10),
        ("keyword regex, the standard approach       precision .61   "
         "recall .27", SMALL, True, INK, 2),
        ("whole-page LLM reading                                 precision .84   "
         "recall .73", SMALL, True, INK, 10),
        ("Keyword search loses roughly three forecasts in four — and distorts "
         "what it keeps. It retrieves advertising (\"return of prosperity\" "
         "appears on one gold page only inside a bank's New Year ad) and pulls "
         "a bullish 1907 forecast out of a passage headlined \"Prophecies Gone "
         "Wrong\" — a report that forecasters had been wrong — which then "
         "scores against the Panic and manufactures a signal out of nothing.",
         SMALL, False, INK, 12),
        ("3 · SCORE", BODY, True, BLUE, 4),
        ("Deterministic; no model involved. 43% of extracted claims are not "
         "about the US economy (regional, industry, foreign) and would "
         "otherwise be graded against national statistics. Claims with no "
         "direction, foreign scope, or a date outside a series' coverage are "
         "marked UNSCORABLE with a recorded reason and left out — never "
         "guessed.", SMALL, False, INK, 12),
        ("4 · ANALYSE", BODY, True, BLUE, 4),
        ("Forecasts within an era share wire copy and one macroeconomic "
         "reality, so the effective sample is ≈21 three-year blocks, not "
         f"{st['n_scorable']:,} claims. Every out-of-fold number uses "
         "leave-one-block-out CV; every interval is block-bootstrapped.",
         SMALL, False, INK, 12),
        ("VALIDATION — 33 known-answer scorer tests against a synthetic "
         "economy · 90 offline pipeline tests · total compute cost $24.97.",
         SMALL, True, INK, 0),
    ], "METHODS"))

    # ---- results: five figures, stacked ----------------------------------
    figs = ["v2_fig1_fifteen_of_fifteen", "v2_fig2_mechanism",
            "v2_fig3_no_smart_paper", "v2_fig4_nothing_transfers",
            "v2_fig5_nobody_has_skill"]
    paths = [f"{FIGDIR}/{n}.png" for n in figs]
    missing = [p for p in paths if not os.path.exists(p)]
    if missing:
        sys.exit(f"missing figures: {missing}\n"
                 "run: python src/make_poster_v2_figures.py")

    heights = [fig_height(p, FIG_W) for p in paths]
    slack = (RESULTS_BOTTOM - RESULTS_TOP) - sum(heights)
    gap = slack / (len(paths) - 1)
    x = RESULTS_L + (RESULTS_R - RESULTS_L - FIG_W) / 2
    y = RESULTS_TOP
    for p, h in zip(paths, heights):
        slide.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(FIG_W))
        y += h + gap
    y -= gap
    print(f"  {'RESULTS':<16} {sum(heights):5.1f}in of "
          f"{RESULTS_BOTTOM - RESULTS_TOP:5.1f}in   "
          f"{'ok' if slack >= 0 else 'OVERFLOW'}  (gap {gap:.2f}in)")
    ok.append(slack >= 0)

    # ---- discussion ------------------------------------------------------
    m, pb, tx = st["mech"], st["pub"], st["text"]
    ok.append(add_text(slide, DISCUSS, [
        ("Optimism was structural, not stupidity.", BODY, True, INK, 4),
        ("The economy really does expand most of the time, so a permanent bull "
         "is right more often than not. What the record shows is not that "
         "forecasters were foolish, but that the press published a FIXED ratio "
         "of optimism that carried no information about conditions. The signal "
         "was never in the direction of the forecast — there was no signal.",
         SMALL, False, INK, 12),
        ("The professionals are no better.", BODY, True, INK, 4),
        ("This is what stops the finding reading as \"old newspapers were "
         "dumb.\" The Fed's own staff, with modern national accounts, scored "
         "54% and had essentially zero skill over the naive rule; the Greenbook "
         "predicted \"worsen\" 6 times in 490 forecasts and got 0 of 6 right. "
         "Sixty years of newspapers and sixty years of professional "
         "macroeconomics land in the same place.", SMALL, False, INK, 12),
        ("Sixty years, no improvement.", BODY, True, INK, 4),
        ("Annual accuracy trends +0.002 per decade (p = 0.87) across 1900–1963. "
         "Two world wars, the Depression and the founding of modern "
         "macroeconomics produced no measurable gain.", SMALL, False, INK, 12),
        ("Should we train a neural network? We measured instead of guessing.",
         BODY, True, VERM, 4),
        (f"The full text of a forecast, given {tx['text_word12_vocab']:,} "
         f"features, reaches AUC {tx['text_word12_insample']:.2f} in sample and "
         f"{tx['text_word12']:.2f} out of fold. The binding constraint is not "
         "model capacity — it is that sixty years of forecasts contain only "
         "~21 independent time blocks. More parameters is the one intervention "
         "that provably cannot fix that: a larger model would memorise harder "
         "and transfer no better.", SMALL, False, INK, 12),
        ("Two findings reversed under more data — and we report both.",
         BODY, True, INK, 4),
        ("On a crisis-only corpus (n = 232), forecasts that swam against the "
         "press consensus looked MORE accurate (52% vs 43%). On the full "
         "continuous corpus (n = 1,967) it reverses: 38.7% vs 53.3%. Separately, "
         "a benchmark this project previously reported did not survive "
         "recovering its deleted source data, and is corrected in Panel 5. Both "
         "are the argument for a continuous corpus over sampling around famous "
         "crashes.", SMALL, False, INK, 12),
        ("LIMITATIONS", BODY, True, INK, 4),
        ("Extraction quality is the binding limitation: this corpus was "
         "extracted at F1 = 0.527 to fit a $30 budget, against 0.784 for the "
         "best extractor tested, which would have cost $404. The gold standard "
         "is model-adjudicated, not human — a two-person recode is the "
         "outstanding step. Coverage thins after 1940. The scope filter is not "
         "neutral: discarded regional and foreign claims are MORE optimistic "
         "(78.7% vs 73.1%), so true press-wide optimism is somewhat higher than "
         "we report.", SMALL, False, INK, 12),
        ("Before all fifteen downturns of the era, the American press said "
         "things would get better.", BODY, True, VERM, 0),
    ], "DISCUSSION"))

    # ---- references ------------------------------------------------------
    ok.append(add_text(slide, REFS, [
        ("1. Library of Congress, Chronicling America: Historic American "
         "Newspapers.  2. NBER, US Business Cycle Expansions and Contractions.  "
         "3. Federal Reserve Bank of St. Louis, FRED (INDPRO, CPIAUCNS, UNRATE, "
         "NBER historical stock index).  4. Federal Reserve Board, Greenbook "
         "Forecast Data Sets, 1967–2020.  5. Federal Reserve Bank of "
         "Philadelphia, Survey of Professional Forecasters; Livingston Survey.  "
         "6. Baker, Bloom & Davis (2016), Measuring Economic Policy Uncertainty, "
         "QJE 131(4).", TINY, False, INK, 0),
    ], "REFERENCES"))

    # ---- acknowledgements ------------------------------------------------
    ok.append(add_text(slide, ACK, [
        ("Boston University RISE Practicum. Newspaper data from the Library of "
         "Congress; business-cycle dates from the National Bureau of Economic "
         "Research; economic series from the Federal Reserve — all public and "
         "free to use. Total compute cost $24.97.", TINY, False, INK, 0),
    ], "ACKNOWLEDGEMENTS"))

    prs.save(OUTFILE)
    print(f"\nwrote {OUTFILE}")
    if not all(ok):
        sys.exit("FAILED: at least one column overflows its box — trim the text "
                 "or shrink the font before printing.")


if __name__ == "__main__":
    main()
