"""
Render a built .pptx poster to a PNG so the layout can actually be looked at.

There is no LibreOffice on this machine, and a poster that overflows a column is
invisible until it is printed -- so this reads the deck's shape geometry directly
and re-draws it: pictures at their real positions, text boxes with their runs
wrapped at the box width and the same point sizes. It is a layout proof, not a
faithful renderer (fonts and kerning differ), which is exactly what is needed to
catch a column running past its box.

Usage:
    python src/preview_poster.py RISE_Poster_2026_v2.pptx
"""

import sys
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

# The template's own section-header boxes ("Introduction", "Results", ...) are
# sized loosely around large display text; this crude wrapper reports them as
# overflowing when they do not. Skip them so a real overflow stands out.
TEMPLATE_HEADERS = {"TextBox 19", "TextBox 21", "TextBox 22", "TextBox 25",
                    "TextBox 26", "TextBox 27"}

PPI = 36  # preview pixels per poster inch -> a 36x48in poster becomes 1296x1728
EMU_IN = 914400

FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
]


def load_font(size_px, bold):
    for p in ([FONT_CANDIDATES[0]] if bold else []) + FONT_CANDIDATES[1:]:
        try:
            return ImageFont.truetype(p, max(6, int(size_px)))
        except OSError:
            continue
    return ImageFont.load_default()


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else "RISE_Poster_2026_v2.pptx"
    prs = Presentation(path)
    W = int(prs.slide_width / EMU_IN * PPI)
    H = int(prs.slide_height / EMU_IN * PPI)
    img = Image.new("RGB", (W, H), "white")
    dr = ImageDraw.Draw(img)
    slide = prs.slides[0]

    def px(emu):
        return int(emu / EMU_IN * PPI)

    for sh in slide.shapes:
        if sh.left is None:
            continue
        L, T = px(sh.left), px(sh.top)
        Wd, Ht = px(sh.width), px(sh.height)

        if sh.shape_type == 13:  # PICTURE
            try:
                im = Image.open(sh.image.blob and __import__("io").BytesIO(
                    sh.image.blob))
                img.paste(im.convert("RGB").resize((Wd, Ht)), (L, T))
            except Exception as e:
                dr.rectangle([L, T, L + Wd, T + Ht], outline="red", width=2)
                dr.text((L + 4, T + 4), f"[image {e}]", fill="red")
            continue

        # Boxes from the template: draw the outline so overflow is visible.
        if sh.shape_type == 1:
            dr.rounded_rectangle([L, T, L + Wd, T + Ht], radius=8,
                                 outline=(150, 150, 150), width=2)
            continue

        if not sh.has_text_frame:
            continue

        y = T
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                size_pt = run.font.size.pt if run.font.size else 18
                size_px = size_pt / 72 * PPI
                font = load_font(size_px, bool(run.font.bold))
                try:
                    color = run.font.color.rgb
                    fill = (color[0], color[1], color[2]) if not isinstance(
                        color, str) else "#" + str(color)
                except Exception:
                    fill = (11, 11, 11)
                if isinstance(fill, tuple) is False:
                    fill = str(fill)
                # The template's section headers set wrap=False with autofit, so
                # they run past their box on one line rather than wrapping.
                if sh.text_frame.word_wrap is False:
                    lines = run.text.split("\n")
                else:
                    approx_cpl = max(6, int(Wd / (size_px * 0.50)))
                    lines = textwrap.wrap(run.text, approx_cpl) or [""]
                centered = para.alignment == PP_ALIGN.CENTER
                for line in lines:
                    if centered:
                        tw = dr.textlength(line, font=font)
                        dr.text((L + (Wd - tw) / 2, y), line, font=font, fill=fill)
                    else:
                        dr.text((L, y), line, font=font, fill=fill)
                    y += size_px * 1.18
            y += (para.space_after.pt if para.space_after else 0) / 72 * PPI
        if y > T + Ht + 2 and sh.name not in TEMPLATE_HEADERS:
            dr.rectangle([L, T, L + Wd, T + Ht], outline="red", width=3)
            dr.text((L, T + Ht), " TEXT OVERFLOWS THIS BOX", fill="red",
                    font=load_font(18, True))

    out = path.replace(".pptx", "_preview.png")
    img.save(out)
    print("wrote", out, img.size)


if __name__ == "__main__":
    main()
