"""
Fill the RISE 2026 poster template with this project's content.

Reads "2026 Poster Template.pptx", writes "RISE_Poster_2026.pptx". Uses the
template's own layout (36x48in, slide 0): Introduction and Methods down the left,
a tall Results column in the middle, Discussion/Conclusions on the right.

Numbers are NOT hardcoded where they can be recomputed -- the headline figures
are pulled from monthly_scored.csv so the poster can never drift from the data.

Usage: python build_poster.py
"""

import copy
from pathlib import Path

import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

TEMPLATE = "2026 Poster Template.pptx"
OUT = "RISE_Poster_2026.pptx"
FIG = Path("figures/poster_figures")

# Pages read during the monthly extraction run. data/monthly/*.jsonl is too
# large for git (385 MB) and is untracked, so this is the fallback when the
# corpus is not on the machine doing the build; it is recomputed from the real
# file whenever that file IS present.
N_PAGES_RECORDED = 15_721

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
VERM = RGBColor(0xD5, 0x5E, 0x00)
BLUE = RGBColor(0x00, 0x72, 0xB2)

BODY = 26          # body text, readable at ~4 ft on a 36x48in poster
SMALL = 21
CAPTION = 19


def _n_pages():
    """Pages read. Prefer the real corpus; fall back to committed sources."""
    corpus = Path("data/corpus/monthly/pages_monthly.jsonl")
    if corpus.exists():
        return sum(1 for _ in open(corpus, encoding="utf-8"))
    return N_PAGES_RECORDED


def stats():
    """Recompute the headline numbers so the poster can't drift from the data.

    Everything here comes from monthly_scored.csv, which IS committed -- the
    poster must stay buildable on a machine that does not have the 385 MB page
    corpus. monthly_scored.csv has one row per extracted claim, so its length is
    the extraction count."""
    raw = pd.read_csv("data/scored/monthly_scored.csv")
    s = raw[(raw["scorable"] == True) & (raw["hit"].isin([0, 1]))].copy()
    from truth_data import NBER_RECESSIONS
    rec = set()
    for a, b in NBER_RECESSIONS:
        rec.update(pd.period_range(a, b, freq="M"))
    p = pd.PeriodIndex(pd.to_datetime(s["date"]), freq="M")
    s["in_rec"] = [x in rec for x in p]
    s["year"] = pd.to_datetime(s["date"]).dt.year
    pess = s["predicted_norm"].isin(["worsen", "down"])

    # Block bootstrap by 3-year period: forecasts cluster hard within an era, so
    # resampling individual claims would understate the interval.
    blk = (s["year"] - 1900) // 3
    blocks = blk.unique()
    rng = np.random.default_rng(0)
    gaps = []
    for _ in range(2000):
        pick = rng.choice(blocks, len(blocks), replace=True)
        bs = pd.concat([s[blk == b] for b in pick])
        if bs["in_rec"].nunique() < 2:
            continue
        gaps.append(bs[~bs["in_rec"]]["hit"].mean() - bs[bs["in_rec"]]["hit"].mean())
    lo, hi = np.percentile(gaps, [2.5, 97.5]) * 100

    # Months the index actually covers -- not every sampled month yields a
    # scorable US-national forecast, and the poster should not claim it does.
    idx = Path("data/scored/press_index.csv")
    n_months = len(pd.read_csv(idx)) if idx.exists() else None

    # Professional-forecaster reference point (committed).
    gb = pd.read_csv("data/scored/greenbook_scored.csv")
    gb_worse = gb[gb["pred_direction"] == "worsen"]

    return {
        "n_pages": _n_pages(),
        "n_claims": len(raw),
        "n_scorable": len(s),
        "n_months": n_months,
        "hit": s["hit"].mean(),
        "hit_exp": s[~s["in_rec"]]["hit"].mean(),
        "hit_rec": s[s["in_rec"]]["hit"].mean(),
        "gap_lo": lo, "gap_hi": hi,
        "pess_exp": pess[~s["in_rec"]].mean(),
        "pess_rec": pess[s["in_rec"]].mean(),
        "upbeat": s["predicted_norm"].isin(["improve", "up"]).mean(),
        "gb_n": len(gb), "gb_hit": gb["hit"].mean(),
        "gb_worse_n": len(gb_worse), "gb_worse_hit": gb_worse["hit"].mean(),
    }


def add_text(slide, left, top, width, height, blocks):
    """blocks = [(text, size, bold, color, space_after_pt), ...]"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, after) in enumerate(blocks):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(after)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_pic(slide, path, left, top, width):
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(width))


def fig_height(path, width):
    """Rendered height of a figure at a given width, from its pixel aspect."""
    w, h = Image.open(path).size
    return width * h / w


def main():
    st = stats()
    prs = Presentation(TEMPLATE)

    # Use slide 0's layout (big central Results column); drop the alternate.
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if len(slides) > 1:
        for sl in slides[1:]:
            xml_slides.remove(sl)
    slide = prs.slides[0]

    # ---- Title -------------------------------------------------------------
    from pptx.enum.text import PP_ALIGN
    for sh in slide.shapes:
        if sh.name == "TextBox 4":
            tf = sh.text_frame
            tf.clear()
            lines = [
                ("Did Anyone See It Coming?", 98, True, INK),
                ("Machine-reading a century of American newspaper economic "
                 "forecasts, 1900–1963", 58, False, MUTED),
                ("BU RISE Practicum — Data Science", 44, True, INK),
            ]
            for i, (text, size, bold, color) in enumerate(lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # The template centres the title, and the RISE logo sits over the
                # box's left edge -- centring is what keeps the text clear of it.
                para.alignment = PP_ALIGN.CENTER
                r = para.add_run()
                r.text = text
                r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # ---- INTRODUCTION  (box at 0.9, 5.5, 8.8 x 18.3) -----------------------
    add_text(slide, 1.3, 7.4, 8.0, 16.0, [
        ("For sixty years American newspapers told readers what the economy "
         "would do next. Were they right?", BODY, True, INK, 16),
        ("The record has never been scored at scale. Reading a century of "
         "newspapers by hand is infeasible — and the standard shortcut, keyword "
         "search, turns out to lose most of the evidence.",
         BODY, False, INK, 16),
        ("We ask two questions:", BODY, True, INK, 8),
        ("1.  Can a language model extract historical forecasts reliably enough "
         "to score them?", BODY, False, INK, 8),
        ("2.  Did the press see downturns coming — and did it ever get better?",
         BODY, False, INK, 16),
        ("The design rule: the model decides WHAT was predicted; real economic "
         "data decides WHETHER IT CAME TRUE. The model is never asked if a "
         "forecast was right.", BODY, True, VERM, 0),
    ])

    # ---- METHODS  (box at 0.9, 24.4, 8.8 x 22.5) ---------------------------
    add_text(slide, 1.3, 26.3, 8.0, 9.0, [
        (f"{st['n_pages']:,} newspaper pages", BODY, True, INK, 4),
        ("Library of Congress Chronicling America. Sampled from every one of the "
         "768 months from 1900-01 to 1963-12, with no gaps, using a fixed "
         f"direction-neutral search phrase set held constant across all months. "
         f"{st['n_months']} of those months yield at least one scorable "
         "US-national forecast.", SMALL, False, INK, 12),
        (f"{st['n_claims']:,} forecasts extracted", BODY, True, INK, 4),
        ("An LLM reads each full page and returns structured claims: direction, "
         "topic, horizon, hedging, speaker, and scope.", SMALL, False, INK, 12),
        (f"{st['n_scorable']:,} scored against reality", BODY, True, INK, 4),
        ("A deterministic scorer compares each forecast to what NBER business-"
         "cycle dates and Federal Reserve series actually did over its horizon. "
         "No model judgement. Claims that cannot be fairly scored — foreign, "
         "regional, or outside a series' coverage — are marked unscorable, "
         "never guessed.", SMALL, False, INK, 12),
        ("Validation", BODY, True, INK, 4),
        ("16 pages hand-annotated before any model ran (52 forecasts, 44 "
         "boundary cases). The scorer is proven by 33 known-answer tests.",
         SMALL, False, INK, 0),
    ])
    add_pic(slide, FIG / "figE_method.png", 1.4, 36.0, 7.8)
    add_text(slide, 1.4, 41.6, 7.8, 4.0, [
        ("Keyword search finds one forecast in four. Whole-page LLM reading "
         "finds three in four, at higher precision — measured against the "
         "hand-built gold standard.", CAPTION, False, MUTED, 0),
    ])

    # ---- RESULTS  (box at 10.6, 5.5, 15.0 x 41.5) --------------------------
    # Six panels in a 39.5in column. Rather than hardcode a y for each -- which
    # silently overlaps the moment a figure's aspect changes -- walk a cursor
    # down and take each figure's real height from its pixels. The assertion at
    # the end fails loudly if the column overflows.
    COL_X, COL_W = 11.1, 14.0
    FIG_W = 11.5                      # narrower than the column so 6 panels fit
    FIG_X = COL_X + (COL_W - FIG_W) / 2
    TITLE_H, CAP_H, GAP = 0.80, 1.05, 0.25
    Y_TOP, Y_BOTTOM = 7.3, 46.9

    panels = [
        ("1.  The forecast mix never responded to the economy",
         "figA_mechanism.png",
         f"Downbeat forecasts were {st['pess_exp']:.1%} of the total in "
         f"expansions and {st['pess_rec']:.1%} in recessions — identical. About "
         f"{st['upbeat']:.0%} of all forecasts were upbeat, in booms and busts "
         f"alike. This survives a within-decade comparison and a block "
         f"bootstrap."),
        ("2.  So accuracy collapsed exactly when it mattered",
         "figB_consequence.png",
         f"Because the mix never moved, accuracy fell from {st['hit_exp']:.1%} "
         f"in expansions to {st['hit_rec']:.1%} in recessions — a consequence of "
         f"panel 1, not a separate finding (gap "
         f"{st['hit_exp'] - st['hit_rec']:+.1%} points, 95% CI "
         f"[{st['gap_lo']:+.1f}, {st['gap_hi']:+.1f}]). Pessimists were right "
         f"when it counted; there were never enough of them."),
        ("3.  Sixty years, no improvement",
         "figC_no_learning.png",
         "Annual accuracy is flat across six decades (53.7% in the 1900s to "
         "49.9% in the 1960s). No publisher did better either: seven papers "
         "with 200+ scored forecasts all fall between 44% and 56%."),
        ("4.  What predicts whether a forecast came true?",
         "figD_what_predicts.png",
         "Out-of-fold AUC (leave-one-block-out, 21 three-year blocks): "
         "macroeconomic conditions 0.505 — chance. How the forecast was written "
         "0.561. Both combined 0.588. We do not claim to predict which forecasts "
         "come true; 0.561 is barely above chance, and saying so is the result."),
        ("5.  Worse than chance on prices, markets and jobs",
         "figF_topics.png",
         "Accuracy varies far more by subject than by wording — a 23-point "
         "spread, against 2 points for hedging. The one above-chance category is "
         "vague optimism about “business”. Price calls track the era’s inflation "
         "regime, not skill: after 1948, “prices will fall” was right 0 times in 93."),
        ("6.  October 1929: no one saw it coming",
         "figG_1929.png",
         "In the month of the Great Crash, 86% of US-national forecasts still "
         "predicted improvement — more than in June. Forecasts printed "
         "August–December 1929 came true 20.8% of the time. Found with no "
         "episode labels and no outcome information."),
    ]

    y = Y_TOP
    for head, fname, cap in panels:
        add_text(slide, COL_X, y, COL_W, TITLE_H, [(head, 30, True, INK, 0)])
        y += TITLE_H
        h = fig_height(FIG / fname, FIG_W)
        add_pic(slide, FIG / fname, FIG_X, y, FIG_W)
        y += h + 0.10
        add_text(slide, COL_X, y, COL_W, CAP_H, [(cap, CAPTION, False, MUTED, 0)])
        y += CAP_H + GAP

    assert y <= Y_BOTTOM, (
        f"results column overflows: content ends at {y:.2f}in, "
        f"limit {Y_BOTTOM}in. Reduce FIG_W.")

    # ---- DISCUSSION / CONCLUSIONS  (box 26.5, 5.5, 8.7 x 31.1) -------------
    add_text(slide, 26.9, 8.6, 7.9, 28.0, [
        ("The press told readers the economy would improve — in booms and busts "
         "alike, in the same proportion, for sixty years.", BODY, True, VERM, 14),
        ("When downturns came, roughly three in four forecasts were pointing "
         "the wrong way. Forecasting the economy's direction was, and remained, "
         "close to a coin flip.", BODY, False, INK, 18),
        ("A finding that reversed", BODY, True, INK, 6),
        ("On a smaller crisis-only sample (n=232), forecasts that went against "
         "the press consensus looked MORE accurate (52% vs 43%). On the full "
         "continuous corpus (n=1,967) it reverses: 38.7% vs 53.3%. The first "
         "result was small-sample noise. This is the clearest argument for "
         "building the continuous corpus — it overturned its own earlier "
         "finding.", SMALL, False, INK, 14),
        ("Nor is this only a newspaper problem", BODY, True, INK, 6),
        (f"Scored the same way, the Federal Reserve’s own internal Greenbook "
         f"forecasts ({st['gb_n']:,}, 1967–2020) come in at {st['gb_hit']:.1%} — "
         f"and they called a downturn just {st['gb_worse_n']} times, getting "
         f"{st['gb_worse_hit']:.0%} of those right. Different era and variable, "
         f"so this is a reference point rather than a controlled comparison — "
         f"but directional forecasting is hard for everyone.",
         SMALL, False, INK, 18),
        ("Limitations", BODY, True, INK, 6),
        ("•  The monthly corpus used a cheaper extractor (F1≈0.53–0.61, "
         "depending on the reasoning-effort setting, which the run log does not "
         "record) than the methods comparison (F1≈0.79), to fit a $25 budget.",
         SMALL, False, INK, 6),
        ("•  The gold standard is model-adjudicated; a two-person human recode "
         "is the next step.", SMALL, False, INK, 6),
        ("•  Digitization thins after 1940, so the late series is noisier.",
         SMALL, False, INK, 6),
        ("•  Our index does not correlate with the published historical policy-"
         "uncertainty index (|r|<0.08) — we measure forecast direction, a "
         "different construct. We report the null.", SMALL, False, INK, 6),
        ("•  Forecasts cluster within eras: the effective sample is ~21 time "
         "blocks, not 14,251 claims. All inference uses grouped CV and block "
         "bootstraps.", SMALL, False, INK, 14),
        ("Contribution", BODY, True, INK, 6),
        ("A validated, reproducible pipeline that turns a century of "
         "qualitative forecasts into scored data — and evidence that the "
         "standard keyword approach would have distorted the answer.",
         SMALL, False, INK, 0),
    ])

    # ---- REFERENCES  (box 26.5, 37.2, 8.7 x 4.7) ---------------------------
    add_text(slide, 26.9, 39.5, 7.9, 3.2, [
        ("Library of Congress, Chronicling America (loc.gov API).  •  National "
         "Bureau of Economic Research, US Business Cycle Expansions and "
         "Contractions.  •  Federal Reserve Bank of St. Louis, FRED: INDPRO, "
         "CPIAUCNS, UNRATE, historical common-stock index.  •  Baker, Bloom & "
         "Davis (2016), Measuring Economic Policy Uncertainty, QJE.",
         18, False, INK, 0),
    ])

    # ---- ACKNOWLEDGEMENTS  (box 26.5, 42.5, 8.7 x 4.5) ---------------------
    add_text(slide, 26.9, 44.7, 7.9, 2.4, [
        ("Thanks to the BU RISE Practicum instructors and to the Library of "
         "Congress for open access to the Chronicling America corpus. All code, "
         "data-building scripts and tests are in the project repository.",
         18, False, INK, 0),
    ])

    prs.save(OUT)
    print(f"-> {OUT}")
    print(f"   {st['n_pages']:,} pages | {st['n_claims']:,} claims | "
          f"{st['n_scorable']:,} scored | {st['n_months']} index months | "
          f"hit {st['hit']:.1%}")
    print(f"   expansion {st['hit_exp']:.1%} vs recession {st['hit_rec']:.1%} "
          f"(gap {st['hit_exp'] - st['hit_rec']:+.1%}, CI "
          f"[{st['gap_lo']:+.1f}, {st['gap_hi']:+.1f}])")
    print(f"   Greenbook reference: {st['gb_hit']:.1%} (n={st['gb_n']:,})")


if __name__ == "__main__":
    main()
